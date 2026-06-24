#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Инвест-дек AIVibe (python-pptx). 16:9, тёмные титулы + светлый контент,
5 параллельных слайдов-сценариев (как работает + мокап) и графики доходности."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.oxml.ns import qn

OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "docs", "AIVibe_Investor_Deck_2026-06.pptx")

# ── Палитра ────────────────────────────────────────────────────────
INK   = RGBColor(0x14, 0x21, 0x3D)   # тёмно-синий (фон титулов)
NAVY  = RGBColor(0x1E, 0x2A, 0x4A)
STEEL = RGBColor(0x2A, 0x6F, 0x97)   # синий акцент
ICE   = RGBColor(0xCA, 0xDC, 0xFC)   # светло-голубой
GOLD  = RGBColor(0xC9, 0xA2, 0x4A)   # золотой акцент
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xF5, 0xF7, 0xFB)   # светлый фон контента
CARD  = RGBColor(0xEE, 0xF2, 0xF8)
LINE  = RGBColor(0xCF, 0xD8, 0xE3)
MUTED = RGBColor(0x5A, 0x64, 0x72)
KEEP  = RGBColor(0x1E, 0x7A, 0x52)   # зелёный
KILL  = RGBColor(0x9C, 0x38, 0x48)   # приглушённый красный
INKTX = RGBColor(0x16, 0x20, 0x33)

HEAD = "Cambria"   # заголовки (safe-list serif)
BODY = "Arial"     # текст (safe-list sans)

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ── Хелперы ────────────────────────────────────────────────────────
def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return r

def box(s, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE, radius=0.10):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp

def text(s, x, y, w, h, runs, size=14, color=INKTX, bold=False, font=BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, leading=1.06, space=2):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in ("left","right","top","bottom"):
        setattr(tf, f"margin_{m}", 0)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = leading; p.space_after = Pt(space)
        if isinstance(para, str):
            para = [(para, {})]
        for seg in para:
            tx, opt = (seg, {}) if isinstance(seg, str) else seg
            r = p.add_run(); r.text = tx
            r.font.name = opt.get("font", font)
            r.font.size = Pt(opt.get("size", size))
            r.font.bold = opt.get("bold", bold)
            r.font.italic = opt.get("italic", False)
            r.font.color.rgb = opt.get("color", color)
    return tb

def kicker(s, x, y, label, color=GOLD):
    text(s, x, y, 8, 0.3, label.upper(), size=11.5, color=color, bold=True, font=BODY)

def chip(s, x, y, label, color, w=2.0, h=0.42):
    c = box(s, x, y, w, h, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(s, x, y-0.01, w, h, label, size=12.5, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c

def numdot(s, x, y, n, d=0.42, fill=GOLD, fg=INK):
    box(s, x, y, d, d, fill=fill, shape=MSO_SHAPE.OVAL)
    text(s, x, y-0.02, d, d, str(n), size=15, color=fg, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD)

def rarrow(s, x, y, w=0.5, h=0.28, color=STEEL):
    a = box(s, x, y, w, h, fill=color, shape=MSO_SHAPE.CHEVRON)
    return a

def content_header(s, kick, title, title_size=33, tcolor=INK):
    bg(s, PAPER)
    kicker(s, 0.7, 0.55, kick)
    text(s, 0.7, 0.84, 12.0, 0.9, title, size=title_size, color=tcolor, bold=True, font=HEAD)

PAGE = [1]  # титул = 1; footer() авто-инкрементит
def footer(s, *_, dark=False):
    PAGE[0] += 1
    col = ICE if dark else MUTED
    text(s, 0.7, 7.08, 8.0, 0.3, "AIVibe · Инвест-презентация · конфиденциально",
         size=8.5, color=col, font=BODY)
    text(s, 11.6, 7.08, 1.05, 0.3, f"{PAGE[0]:02d}", size=8.5, color=col, font=BODY,
         align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════
# 01 — ТИТУЛ
# ════════════════════════════════════════════════════════════════════
s = slide(); bg(s, INK)
box(s, 0.9, 2.05, 0.9, 0.12, fill=GOLD)
text(s, 0.9, 2.35, 11, 1.4, "AIVibe", size=66, color=WHITE, bold=True, font=HEAD)
text(s, 0.95, 3.7, 11.5, 1.0,
     "AI-дизайнер пространства для рынка РФ: рендер, расстановка по эргономике и AR —"
     " как инструмент профессионала, а не игрушка для разового покупателя.",
     size=17, color=ICE, font=BODY, leading=1.2)
chip(s, 0.95, 5.15, "Стратегия выхода на рынок", STEEL, w=3.3, h=0.5)
chip(s, 4.4, 5.15, "5 сценариев · экономика", NAVY, w=3.3, h=0.5)
text(s, 0.95, 6.55, 11, 0.4, "Инвест-презентация · июнь 2026 · КОНФИДЕНЦИАЛЬНО",
     size=11.5, color=RGBColor(0x9F,0xAD,0xC2), bold=True)

# ════════════════════════════════════════════════════════════════════
# 02 — РЫНОК: большой, но «кладбище»
# ════════════════════════════════════════════════════════════════════
s = slide(); content_header(s, "Рынок и проблема", "Большой спрос — и высокая смертность ниши")
# три стат-callout
stats = [("~700 млрд ₽", "рынок мебели РФ, 2024 (+20% г/г)"),
         ("×4,7", "рост рынка AI-interior к 2032 ($1,47→6,96 млрд)"),
         ("раз в 5–10 лет", "частота покупки мебели — корень всех смертей")]
cx = 0.7
for big, small in stats:
    box(s, cx, 1.95, 3.85, 1.5, fill=WHITE, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    text(s, cx+0.25, 2.12, 3.4, 0.7, big, size=27, color=STEEL, bold=True, font=HEAD)
    text(s, cx+0.25, 2.78, 3.4, 0.55, small, size=12, color=MUTED)
    cx += 4.06
text(s, 0.7, 3.75, 11.9, 0.5,
     [[("Та же ниша «AR/AI-дизайн интерьера» — кладбище: на ней по ", {}),
       ("одной и той же причине", {"bold": True, "color": INK}),
       (" легли проекты с любым размером финансирования.", {})]],
     size=14.5, color=INKTX)
# полоски сожжённого венчура
dead = [("Modsy", "$72,7M сожжено · закрылась 2022"),
        ("Hutch", "не вытянул комиссию → ушёл в игры"),
        ("Havenly", "выжила, лишь став ритейл-холдингом"),
        ("Faradise (РФ)", "AR-маркетплейс умер ~2019 → B2B-студия")]
dy = 4.45
for name, fate in dead:
    box(s, 0.7, dy, 0.16, 0.42, fill=KILL)
    text(s, 1.0, dy-0.02, 3.0, 0.42, name, size=13.5, color=INK, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.0, dy-0.02, 8.4, 0.42, fate, size=12.5, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    dy += 0.56
footer(s, 2)

# ════════════════════════════════════════════════════════════════════
# 03 — ПОЧЕМУ УМИРАЮТ (5 причин)
# ════════════════════════════════════════════════════════════════════
s = slide(); content_header(s, "Почему ниша убивает", "Пять структурных причин смерти")
killers = [
    ("Разовый спрос < стоимости привлечения", "повторных покупок мебели лишь 14,7% → экономика массового пользователя убыточна by design"),
    ("Комиссия с редкой покупки не масштабируется", "не помогли даже $72M: транзакций мало, атрибуция утекает, маржа тонкая"),
    ("AR — паритетная фича, а не защита", "маркетплейсы раздают AR бесплатно; «летающие диваны» отпугивают пользователя"),
    ("Пустой каталог + оцифровка дороже рынка", "двусторонняя платформа умирает с пустой стороной предложения"),
    ("Управление вслепую без аналитики/identity", "нельзя отличить прибыльный сегмент от убыточного → деньги в трубу"),
]
yy = 1.95
for i, (h, d) in enumerate(killers, 1):
    numdot(s, 0.7, yy+0.03, i, d=0.5, fill=GOLD)
    text(s, 1.45, yy, 11.1, 0.4, h, size=15.5, color=INK, bold=True)
    text(s, 1.45, yy+0.4, 11.1, 0.4, d, size=12.5, color=MUTED)
    yy += 0.97
footer(s, 3)

# ════════════════════════════════════════════════════════════════════
# 04 — ИНСАЙТ / РЕШЕНИЕ (тёмный)
# ════════════════════════════════════════════════════════════════════
s = slide(); bg(s, INK)
kicker(s, 0.9, 0.7, "Ключевой инсайт")
text(s, 0.9, 1.05, 11.5, 1.6,
     [[("Сменить ", {"color": WHITE}), ("плательщика", {"color": GOLD}),
       (" — не чинить воронку.", {"color": WHITE})]],
     size=38, bold=True, font=HEAD)
# две колонки: масса (убыток) vs профессионал
box(s, 0.9, 2.7, 5.5, 3.4, fill=NAVY, line=KILL, lw=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
text(s, 1.25, 2.95, 4.9, 0.5, "Массовый покупатель", size=18, color=WHITE, bold=True, font=HEAD)
text(s, 1.25, 3.5, 4.9, 2.4,
     [["• спрос раз в 5–10 лет"],
      ["• LTV / CAC = 0,03–0,38 — убыток"],
      ["• платная реклама запрещена арифметикой"],
      ["• окупаемость: никогда"]],
     size=14.5, color=ICE, leading=1.5, space=6)
box(s, 6.9, 2.7, 5.5, 3.4, fill=NAVY, line=KEEP, lw=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
text(s, 7.25, 2.95, 4.9, 0.5, "Профессионал с повторяющимся спросом", size=18, color=WHITE, bold=True, font=HEAD)
text(s, 7.25, 3.5, 4.9, 2.4,
     [["• десятки объектов / проектов в год"],
      ["• LTV / CAC = 3–8 — здоровый"],
      ["• платная реклама окупается"],
      ["• окупаемость: 1–6 месяцев"]],
     size=14.5, color=ICE, leading=1.5, space=6)
text(s, 0.9, 6.45, 11.5, 0.5,
     "Мы зарабатываем на риелторе, застройщике и дизайнере — на той же кодовой базе.",
     size=14, color=GOLD, bold=True)
footer(s, 4, dark=True)

# ════════════════════════════════════════════════════════════════════
# 05 — ПРОДУКТ (как устроен core)
# ════════════════════════════════════════════════════════════════════
s = slide(); content_header(s, "Продукт", "Один продукт: вход рендером → результат в бюджете")
flow = [("Фото / запрос", "пользователь даёт фото комнаты или бриф"),
        ("AI-рендер", "«фото → стиль» за секунды (YandexART)"),
        ("Диалог-дизайнер", "уточняет стиль и бюджет (YandexGPT)"),
        ("Расстановка по нормам", "детерминированный движок, без «летающих диванов»"),
        ("Список покупок / AR", "реальные товары в бюджете")]
n = len(flow); cw = 2.18; gap = 0.27; x0 = 0.7; y0 = 2.5
for i, (h, d) in enumerate(flow):
    x = x0 + i*(cw+gap)
    box(s, x, y0, cw, 2.5, fill=WHITE, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
    numdot(s, x+0.18, y0+0.2, i+1, d=0.46, fill=STEEL, fg=WHITE)
    text(s, x+0.18, y0+0.82, cw-0.36, 0.7, h, size=14, color=INK, bold=True, font=HEAD)
    text(s, x+0.18, y0+1.5, cw-0.36, 0.9, d, size=11, color=MUTED, leading=1.1)
    if i < n-1:
        rarrow(s, x+cw-0.02, y0+1.05, w=0.31, h=0.34, color=GOLD)
text(s, 0.7, 5.55, 11.9, 0.8,
     [[("Наш ров — не AR (это паритет), а ", {}),
       ("детерминированный движок расстановки", {"bold": True, "color": INK}),
       (" по эргономике: он один соединяет «красивую картинку» с реальной купляемой комнатой.", {})]],
     size=14, color=INKTX, leading=1.25)
footer(s, 5)

# ════════════════════════════════════════════════════════════════════
# 06 — НАШ EDGE (3 карточки)
# ════════════════════════════════════════════════════════════════════
s = slide(); content_header(s, "Несправедливое преимущество", "Три актива, которых нет у других")
edges = [
    ("Движок-ров", "Детерминированная расстановка (Kán & Kaufmann). Считает геометрией, не нейросетью → не копируется «лучшим промптом» и не зависит от цены AI-токенов."),
    ("RU-резидентность", "Западные лидеры (InteriorAI, Spacely) в РФ платно недоступны. Мы принимаем рубли (ЮKassa), хостимся в Yandex Cloud — конкурент не может войти платёжно."),
    ("Живой продукт", "5 облачных функций уже в проде; диалоговый AI на отечественном стеке YandexGPT → GigaChat → CoreML с резервированием (Triplex Fallback)."),
]
cx = 0.7
for h, d in edges:
    box(s, cx, 2.1, 3.85, 3.7, fill=WHITE, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    box(s, cx+0.3, 2.4, 0.7, 0.7, fill=GOLD, shape=MSO_SHAPE.OVAL)
    cx += 4.06
# подписи поверх (отдельно для контроля)
cx = 0.7
for h, d in edges:
    text(s, cx+0.3, 3.3, 3.25, 0.5, h, size=18, color=INK, bold=True, font=HEAD)
    text(s, cx+0.3, 3.95, 3.25, 1.7, d, size=12.5, color=MUTED, leading=1.25)
    cx += 4.06
footer(s, 6)

# ════════════════════════════════════════════════════════════════════
# 07 — 5 ПУТЕЙ (матрица)
# ════════════════════════════════════════════════════════════════════
s = slide(); content_header(s, "Развилка стратегии", "Пять путей — мы оценили каждый")
rows = [
    ("Путь", "Маржа", "Ров", "Скорость", "Решение", None),
    ("S1 · Чистый AI-рендер по подписке", "высокая", "низкий", "высокая", "Не выбран", KILL),
    ("S2 · White-label инструмент фабрикам", "высокая", "средн.", "низкая", "Не выбран", KILL),
    ("S3 · SaaS-инструмент для дизайнеров", "средн.", "средн.", "средн.", "Выбран · второй", KEEP),
    ("S4 · Виртуальный стейджинг недвижимости", "высокая", "средн.", "высокая", "Выбран · основной", KEEP),
    ("S5 · Open-core / API-платформа", "высокая", "высокий", "низкая", "Не выбран", KILL),
]
colx = [0.7, 5.5, 6.9, 8.2, 9.7]; colw = [4.7, 1.4, 1.3, 1.5, 2.9]
ry = 2.15; rh = 0.62
for ri, row in enumerate(rows):
    head = ri == 0
    fill = INK if head else (CARD if ri % 2 else WHITE)
    if ri == 4:
        fill = RGBColor(0xE4, 0xEF, 0xEA)
    box(s, 0.7, ry, 11.9, rh, fill=fill, line=LINE, lw=0.75)
    for ci in range(5):
        val = row[ci]
        if head:
            col = WHITE; bold = True; al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        elif ci == 4:
            col = row[5]; bold = True; al = PP_ALIGN.CENTER
        elif ci == 0:
            col = INK; bold = (ri == 4); al = PP_ALIGN.LEFT
        else:
            col = MUTED; bold = False; al = PP_ALIGN.CENTER
        text(s, colx[ci]+0.12, ry, colw[ci]-0.2, rh, val, size=12.5 if not head else 12,
             color=col, bold=bold, align=al, anchor=MSO_ANCHOR.MIDDLE)
    ry += rh
text(s, 0.7, ry+0.12, 11.9, 0.7,
     "S1 — единственное преимущество принадлежит будущему конкуренту (Яндекс/Сбер). "
     "S2 — инструменту нет платящего спроса. S5 — рынка такого API в РФ пока нет.",
     size=11.5, color=MUTED, leading=1.2)
footer(s, 7)

# ════════════════════════════════════════════════════════════════════
# 08–12 — ПЯТЬ СЦЕНАРИЕВ (как работает + мокап)
# ════════════════════════════════════════════════════════════════════
def phone(s, x, y, w=2.5, h=4.4, top="ФОТО", bottom="РЕНДЕР", tc=CARD, bc=STEEL):
    box(s, x, y, w, h, fill=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    box(s, x+0.12, y+0.18, w-0.24, h-0.36, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    half = (h-0.36)/2 - 0.06
    box(s, x+0.26, y+0.32, w-0.52, half, fill=tc, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    text(s, x+0.26, y+0.32, w-0.52, half, top, size=12, color=MUTED, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD)
    box(s, x+0.26, y+0.42+half, w-0.52, half, fill=bc, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    text(s, x+0.26, y+0.42+half, w-0.52, half, bottom, size=12, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD)

def browser(s, x, y, w=4.6, h=3.2, url="фабрика.рф/диван", btn="Примерить в комнате"):
    box(s, x, y, w, h, fill=WHITE, line=LINE, lw=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    box(s, x, y, w, 0.5, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    for i, c in enumerate([KILL, GOLD, KEEP]):
        box(s, x+0.18+i*0.22, y+0.17, 0.16, 0.16, fill=c, shape=MSO_SHAPE.OVAL)
    box(s, x+1.0, y+0.12, w-1.3, 0.26, fill=WHITE, line=LINE, lw=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(s, x+1.15, y+0.12, w-1.5, 0.26, url, size=9.5, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    box(s, x+0.3, y+0.75, w-2.0, h-1.1, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    text(s, x+0.3, y+0.75, w-2.0, h-1.1, "3D-модель\nтовара", size=12, color=MUTED, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD)
    box(s, x+w-1.55, y+1.0, 1.3, 0.55, fill=STEEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    text(s, x+w-1.55, y+1.0, 1.3, 0.55, btn, size=10, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, x+w-1.55, y+1.7, 1.3, 0.5, fill=RGBColor(0xE4,0xEF,0xEA), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2)
    text(s, x+w-1.55, y+1.7, 1.3, 0.5, "+конверсия", size=10, color=KEEP, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def two_photos(s, x, y, w=4.7, h=3.4, a="ПУСТАЯ\nКОМНАТА", b="ОБСТАВЛЕНО"):
    cw = (w-0.4)/2
    box(s, x, y+0.4, cw, 2.0, fill=CARD, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    text(s, x, y+0.4, cw, 2.0, a, size=12, color=MUTED, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD)
    rarrow(s, x+cw+0.04, y+1.2, w=0.32, h=0.36, color=GOLD)
    box(s, x+cw+0.4, y+0.4, cw, 2.0, fill=STEEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    text(s, x+cw+0.4, y+0.4, cw, 2.0, b, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD)
    box(s, x+cw-0.2, y+2.6, cw+0.4, 0.62, fill=RGBColor(0xE4,0xEF,0xEA), line=KEEP, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2)
    text(s, x+cw-0.2, y+2.6, cw+0.4, 0.62, "→ в объявление Авито / Циан", size=11, color=KEEP, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def dashboard(s, x, y, w=4.7, h=3.5):
    box(s, x, y, w, h, fill=WHITE, line=LINE, lw=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    box(s, x+0.15, y+0.15, 1.35, h-0.3, fill=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    text(s, x+0.28, y+0.35, 1.1, 0.3, "Проекты", size=10.5, color=WHITE, bold=True)
    for i, nm in enumerate(["Клиент А", "Клиент Б", "Клиент В"]):
        fill = GOLD if i == 0 else NAVY
        box(s, x+0.28, y+0.75+i*0.5, 1.08, 0.38, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2)
        text(s, x+0.28, y+0.75+i*0.5, 1.08, 0.38, nm, size=9, color=(INK if i==0 else ICE), bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, x+1.65, y+0.25, w-1.85, 2.0, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    text(s, x+1.65, y+0.25, w-1.85, 2.0, "Рендер + расстановка", size=12, color=MUTED, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD)
    for i, (lbl, c) in enumerate([("AR клиенту", STEEL), ("Смета / покупки", KEEP)]):
        bx = x+1.65 + i*((w-1.85)/2)
        box(s, bx, y+2.45, (w-1.85)/2-0.12, 0.55, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.25)
        text(s, bx, y+2.45, (w-1.85)/2-0.12, 0.55, lbl, size=10.5, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def api_diagram(s, x, y, w=4.7, h=3.5):
    cxw = 1.9
    box(s, x+(w-cxw)/2, y+1.35, cxw, 0.9, fill=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    text(s, x+(w-cxw)/2, y+1.35, cxw, 0.9, "/arrange\nAPI", size=13, color=GOLD, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD)
    nodes = ["AR-приложения", "Фабрики", "Дизайн-тулы"]
    for i, nm in enumerate(nodes):
        nx = x + i*(w/3) + 0.1
        box(s, nx, y+0.1, w/3-0.3, 0.7, fill=WHITE, line=STEEL, lw=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
        text(s, nx, y+0.1, w/3-0.3, 0.7, nm, size=10.5, color=INK, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, x+0.3, y+2.7, w-0.6, 0.6, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    text(s, x+0.3, y+2.7, w-0.6, 0.6, "open-source ядро · усиление через сообщество", size=10.5,
         color=MUTED, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def scenario_slide(page, tag, title, steps, verdict, vcolor, oneliner, mockup, primary=False):
    s = slide(); bg(s, PAPER)
    kicker(s, 0.7, 0.5, f"Сценарий {tag}")
    text(s, 0.7, 0.8, 8.3, 0.9, title, size=26, color=INK, bold=True, font=HEAD)
    chip(s, 0.7, 1.72, verdict, vcolor, w=3.0, h=0.46)
    if primary:
        box(s, 3.85, 1.72, 1.5, 0.46, fill=GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        text(s, 3.85, 1.71, 1.5, 0.46, "★ ТАРАН", size=12, color=INK, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.7, 2.55, 6.2, 0.35, "КАК ЭТО РАБОТАЕТ", size=12, color=STEEL, bold=True)
    yy = 3.0
    for i, (h, d) in enumerate(steps, 1):
        numdot(s, 0.7, yy, i, d=0.44, fill=STEEL, fg=WHITE)
        text(s, 1.35, yy-0.04, 5.6, 0.35, h, size=14, color=INK, bold=True)
        text(s, 1.35, yy+0.32, 5.6, 0.4, d, size=11, color=MUTED, leading=1.05)
        yy += 0.82
    box(s, 0.7, 6.35, 6.2, 0.7, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    text(s, 0.95, 6.35, 5.7, 0.7,
         [[("Почему так: ", {"bold": True, "color": INK}), (oneliner, {"color": MUTED})]],
         size=11.5, anchor=MSO_ANCHOR.MIDDLE, leading=1.1)
    mockup(s)
    footer(s, page)
    return s

# S1
scenario_slide(8, "S1", "Чистый AI-рендер по подписке",
    [("Фото комнаты", "пользователь снимает свою комнату"),
     ("Выбор стиля", "лофт / скандинавия / неоклассика…"),
     ("AI-рендер за 20 сек", "«фото → готовый интерьер»"),
     ("Поделиться «было / стало»", "вирусный бесплатный канал")],
    "Не выбран", KILL,
    "ров принадлежит будущему конкуренту (Яндекс/Сбер) — копируется за один спринт.",
    lambda s: phone(s, 9.4, 1.7, 2.9, 4.9, top="ФОТО", bottom="РЕНДЕР"))

# S2
scenario_slide(9, "S2", "White-label инструмент для фабрик и ритейла",
    [("Каталог фабрики", "товары на сайте фабрики/магазина"),
     ("Виджет AIVibe", "встраивается строкой кода"),
     ("«Примерить + расстановка»", "AR + расчёт по эргономике"),
     ("Рост конверсии карточки", "фабрика платит за инструмент")],
    "Не выбран", KILL,
    "инструменту нет платящего спроса; дешёвый аналог (Arigami) уже от 990 ₽.",
    lambda s: browser(s, 7.7, 2.6, 4.9, 3.4))

# S3
scenario_slide(10, "S3", "Вертикальный SaaS для дизайнеров интерьера",
    [("Бриф клиента", "дизайнер заводит проект"),
     ("AI-рендер + расстановка", "вариант за минуту, без «летающих диванов»"),
     ("AR-презентация клиенту", "показывает с телефона"),
     ("Смета и список покупок", "закрывает сделку быстрее")],
    "Выбран · второй", KEEP,
    "дизайнер платит сам, бесплатно приводит клиентов и гонит заказы фабрикам.",
    lambda s: dashboard(s, 7.7, 2.6, 4.9, 3.5))

# S4 (основной)
scenario_slide(11, "S4", "Виртуальный стейджинг недвижимости",
    [("Фото пустой комнаты", "риелтор снимает с телефона"),
     ("Выбор стиля", "под продажу / под аренду"),
     ("AI-меблировка за 90 сек", "движок ставит мебель по нормам"),
     ("В объявление", "Авито / Циан / Домклик")],
    "Выбран · основной", KEEP,
    "платит риелтор/застройщик с десятками объектов в год — спрос повторяющийся.",
    lambda s: two_photos(s, 7.7, 2.7, 4.7, 3.4), primary=True)

# S5
scenario_slide(12, "S5", "Open-core / API-платформа",
    [("Open-source ядро", "движок расстановки на GitHub"),
     ("API-ключ", "разработчик регистрируется"),
     ("Вызов /arrange", "получает координаты без коллизий"),
     ("Интеграция в продукт", "платит за хостинг/SLA/ассеты")],
    "Не выбран", KILL,
    "рынка такого API в РФ пока нет — категорию пришлось бы создавать годами.",
    lambda s: api_diagram(s, 7.7, 2.7, 4.7, 3.5))

# ════════════════════════════════════════════════════════════════════
# 13 — ВЫБРАННАЯ СТРАТЕГИЯ
# ════════════════════════════════════════════════════════════════════
s = slide(); content_header(s, "Решение", "Стейджинг-таран + дизайнеры на одной кодовой базе")
left = [("S4 · Виртуальный стейджинг — таран", "Лучшая экономика и скорость. Риелторы и застройщики (пакет на ЖК — несущий чек). AR и сканер уходят с критического пути.", GOLD),
        ("S3 · Дизайнеры — второй сегмент", "Тот же движок и рендер. Дизайнер платит сам и бесплатно приводит клиентов → включает партнёрский слой «чужими руками».", STEEL)]
yy = 2.1
for h, d, c in left:
    box(s, 0.7, yy, 7.2, 1.55, fill=WHITE, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    box(s, 1.0, yy+0.22, 0.34, 0.34, fill=c, shape=MSO_SHAPE.OVAL)
    text(s, 1.5, yy+0.18, 6.2, 0.5, h, size=16, color=INK, bold=True, font=HEAD)
    text(s, 1.05, yy+0.72, 6.55, 0.78, d, size=12, color=MUTED, leading=1.2)
    yy += 1.75
# правая колонка: общий стек
box(s, 8.2, 2.1, 4.4, 3.2, fill=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
text(s, 8.5, 2.35, 3.9, 0.5, "Одна кодовая база", size=16, color=GOLD, bold=True, font=HEAD)
text(s, 8.5, 2.95, 3.9, 2.3,
     [["✓ движок расстановки (ров)"],
      ["✓ AI-рендер на YandexART"],
      ["✓ диалоговый AI · Triplex Fallback"],
      ["✓ контур оплат ЮKassa"],
      ["✓ ~70% кода переиспользуется"]],
     size=13, color=ICE, leading=1.55, space=5)
text(s, 0.7, 5.7, 11.9, 0.7,
     "Заостряет более ранний вывод (дизайнеры-first): у стейджинга экономика строго лучше — "
     "быстрее до денег и выше маржа, поэтому он становится денежным тараном.",
     size=12.5, color=MUTED, leading=1.2)
footer(s, 13)

# ════════════════════════════════════════════════════════════════════
# NEW — SaaS МОДУЛИ
# ════════════════════════════════════════════════════════════════════
s = slide(); content_header(s, "Продукт детально", "SaaS для дизайнеров: сквозная цепочка в одном окне")
chain = [("0","Бриф"),("1","Стиль"),("2","Расстановка"),("3","Презентация"),("4","Спец-ция"),("5","Смета"),("6","Закупка")]
n=len(chain); cw=1.55; gap=0.18; x0=0.7; y0=2.15
for i,(num,nm) in enumerate(chain):
    x=x0+i*(cw+gap)
    box(s, x, y0, cw, 0.95, fill=WHITE, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    numdot(s, x+0.12, y0+0.13, num, d=0.34, fill=STEEL, fg=WHITE)
    text(s, x+0.05, y0+0.5, cw-0.1, 0.4, nm, size=11, color=INK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i<n-1: rarrow(s, x+cw-0.03, y0+0.33, w=0.2, h=0.3, color=GOLD)
diffs=[("Расстановка по нормам","детерминированный движок, без «летающих диванов» — повторяемо и редактируемо"),
       ("Смета в ₽ с живыми ценами","спецификация и смета связаны с рендером; цены/наличие РФ-фабрик"),
       ("Один поток без Excel","рендер → расстановка → спецификация → смета → закупка без дублирования")]
cx=0.7
for h,d in diffs:
    box(s, cx, 3.5, 3.85, 2.25, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    box(s, cx+0.3, 3.8, 0.5, 0.5, fill=GOLD, shape=MSO_SHAPE.OVAL)
    text(s, cx+0.3, 4.45, 3.25, 0.5, h, size=14.5, color=INK, bold=True, font=HEAD)
    text(s, cx+0.3, 4.98, 3.25, 0.72, d, size=11.5, color=MUTED, leading=1.2)
    cx+=4.06
text(s, 0.7, 6.05, 11.9, 0.4, "Чего нет ни у Foyr, ни у Spacely, ни у Planoplan — единая цепочка и смета в рублях с ценами РФ-фабрик.", size=12, color=STEEL, bold=True)
footer(s)

# ════════════════════════════════════════════════════════════════════
# NEW — ТЕХПАЙПЛАЙН СТЕЙДЖИНГА
# ════════════════════════════════════════════════════════════════════
s = slide(); content_header(s, "Технология", "Виртуальный стейджинг: как это работает")
steps=[("Фото","комнаты"),("Глубина","+ линии"),("Маски","пол/окна"),("План движка","расстановка"),("Стиль","пресет/реф."),("Генерация","SDXL+ControlNet"),("Апскейл","+ авто-QA")]
n=len(steps); cw=1.55; gap=0.18; x0=0.7; y0=2.05
for i,(h,d) in enumerate(steps):
    x=x0+i*(cw+gap)
    hot=i==3
    box(s, x, y0, cw, 1.2, fill=(RGBColor(0xE4,0xEF,0xEA) if hot else WHITE), line=(KEEP if hot else LINE), lw=(1.4 if hot else 1.0), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    numdot(s, x+0.12, y0+0.12, i+1, d=0.32, fill=(KEEP if hot else STEEL), fg=WHITE)
    text(s, x+0.06, y0+0.5, cw-0.12, 0.36, h, size=10.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    text(s, x+0.06, y0+0.85, cw-0.12, 0.3, d, size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    if i<n-1: rarrow(s, x+cw-0.03, y0+0.44, w=0.2, h=0.3, color=GOLD)
box(s, 0.7, 3.7, 5.85, 2.05, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
text(s, 1.0, 3.95, 5.3, 0.45, "Стиль задаёт пользователь", size=15, color=INK, bold=True, font=HEAD)
text(s, 1.0, 4.45, 5.3, 1.2, "Пресет / текст (MVP) · референс-фото через IP-Adapter (премиум). «Переодевает» комнату диффузия SDXL; YandexGPT только собирает промпт — картинку не рисует.", size=12, color=MUTED, leading=1.25)
box(s, 6.75, 3.7, 5.85, 2.05, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
text(s, 7.05, 3.95, 5.3, 0.45, "Движок решает ГДЕ — без «летающих диванов»", size=15, color=INK, bold=True, font=HEAD)
text(s, 7.05, 4.45, 5.3, 1.2, "Детерминированный ArrangementEngine считает расстановку по нормам и подаёт её как жёсткие маски диффузии. У конкурентов мебель ставит сама нейросеть «на удачу».", size=12, color=MUTED, leading=1.25)
box(s, 0.7, 5.95, 11.9, 0.75, fill=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
text(s, 1.0, 5.95, 11.3, 0.75, [[("Движок: ",{"bold":True,"color":GOLD}),("self-hosted SDXL + ControlNet на Yandex GPU (РФ-резидентно, 152-ФЗ) — managed-API в РФ фото не редактирует. Себестоимость ~7–11 ₽/кадр.",{"color":ICE})]], size=12.5, anchor=MSO_ANCHOR.MIDDLE, leading=1.15)
footer(s)

# ════════════════════════════════════════════════════════════════════
# 14 — ЮНИТ-ЭКОНОМИКА (chart LTV/CAC)
# ════════════════════════════════════════════════════════════════════
s = slide(); content_header(s, "Юнит-экономика", "Почему это сходится: LTV / CAC по сегментам")
cd = CategoryChartData()
cd.categories = ["Массовый\nпокупатель", "Дизайнер", "Риелтор /\nзастройщик"]
cd.add_series("LTV / CAC", (0.3, 5.0, 6.0))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(2.1), Inches(6.6), Inches(4.3), cd)
ch = gf.chart; ch.has_legend = False; ch.has_title = False
plot = ch.plots[0]; plot.gap_width = 90; plot.has_data_labels = True
plot.data_labels.number_format = '0.0'; plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(13); plot.data_labels.font.bold = True; plot.data_labels.font.color.rgb = INK
ser = plot.series[0]
for idx, pt in enumerate(ser.points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = KILL if idx == 0 else KEEP
cat_ax = ch.category_axis; cat_ax.tick_labels.font.size = Pt(11); cat_ax.tick_labels.font.bold = True
cat_ax.tick_labels.font.color.rgb = INK
val_ax = ch.value_axis; val_ax.minimum_scale = 0; val_ax.maximum_scale = 8
val_ax.has_major_gridlines = True; val_ax.tick_labels.font.size = Pt(10); val_ax.tick_labels.font.color.rgb = MUTED
# правая колонка пояснений
box(s, 7.8, 2.1, 4.8, 4.3, fill=WHITE, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
text(s, 8.1, 2.4, 4.2, 0.5, "Норма здоровья: LTV / CAC ≥ 3", size=15, color=INK, bold=True, font=HEAD)
text(s, 8.1, 3.1, 4.2, 3.0,
     [[("Масса", {"bold": True, "color": KILL}), (" — 0,03–0,38: убыток by design,", {})],
      [("платная реклама запрещена арифметикой.", {})],
      [("",{})],
      [("Дизайнер / риелтор", {"bold": True, "color": KEEP}), (" — 3–8:", {})],
      [("окупаемость 1–6 мес, маржа стейджинга ~95%", {})],
      [("и устойчива к росту цен на AI (движок —", {})],
      [("вычисление, а не запрос к нейросети).", {})]],
     size=12.5, color=INKTX, leading=1.25)
text(s, 0.7, 6.6, 11.9, 0.3, "Значения — середины диапазонов из разбора; иллюстрация, не обещание.",
     size=10, color=MUTED)
footer(s, 14)

# ════════════════════════════════════════════════════════════════════
# 15 — ГРАФИК ДОХОДНОСТИ (revenue vs cost, break-even)
# ════════════════════════════════════════════════════════════════════
s = slide(); content_header(s, "Графики доходности", "Реальная модель: выход на безубыточность")
months = [f"М{m}" for m in range(1, 25)]
revenue = [0,0,0,0,99,186,294,423,599,794,1010,1240,1489,1754,2033,2332,2646,2979,3328,3691,4069,4461,4869,5291]
costs   = [626,626,626,626,1136,1148,1165,1186,1211,1700,1734,1772,1815,1862,1912,1968,2028,2092,2160,2233,2309,2390,2474,2562]
cd2 = CategoryChartData(); cd2.categories = months
cd2.add_series("Выручка, тыс ₽/мес", revenue)
cd2.add_series("Затраты, тыс ₽/мес", costs)
gf2 = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.7), Inches(2.05), Inches(8.4), Inches(4.4), cd2)
ch2 = gf2.chart; ch2.has_title = False
ch2.has_legend = True; ch2.legend.position = XL_LEGEND_POSITION.TOP; ch2.legend.include_in_layout = False
ch2.legend.font.size = Pt(11); ch2.legend.font.bold = True
s1, s2 = ch2.plots[0].series
s1.format.line.color.rgb = KEEP; s1.format.line.width = Pt(2.75)
s2.format.line.color.rgb = KILL; s2.format.line.width = Pt(2.5)
ca = ch2.category_axis; ca.tick_labels.font.size = Pt(7); ca.tick_labels.font.color.rgb = MUTED
va = ch2.value_axis; va.tick_labels.font.size = Pt(9); va.tick_labels.font.color.rgb = MUTED
va.has_major_gridlines = True; va.minimum_scale = 0
# вынос — безубыточность
box(s, 9.4, 2.3, 3.2, 1.85, fill=RGBColor(0xE4,0xEF,0xEA), line=KEEP, lw=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
text(s, 9.65, 2.5, 2.75, 0.5, "Безубыточность", size=14, color=KEEP, bold=True, font=HEAD)
text(s, 9.65, 3.0, 2.75, 1.1,
     [[("операционная ", {"color": INK}), ("≈ М15", {"bold": True, "color": KEEP}), (";", {"color": INK})],
      [("полный возврат вложенного ", {"color": INK}), ("≈ М23", {"bold": True, "color": KEEP}), (".", {"color": INK})]],
     size=11.5, leading=1.25)
box(s, 9.4, 4.3, 3.2, 2.0, fill=WHITE, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
text(s, 9.65, 4.5, 2.75, 1.7,
     [[("Капитал до окупаемости:", {"bold": True, "color": INK})],
      [("дно «ямы» ≈ 9,3 млн ₽", {"color": MUTED})],
      [("(консерватив ~16 млн).", {"color": MUTED})],
      [("Подушка 18–20 млн ₽.", {"bold": True, "color": STEEL})]],
     size=11, leading=1.3)
text(s, 0.7, 6.62, 11.9, 0.3,
     "Модель на официальном прайсе Yandex Cloud + зарплатах РФ 2025; консервативные допущения, не финансовое обещание.",
     size=10, color=MUTED)
footer(s)

# ════════════════════════════════════════════════════════════════════
# NEW — КАПИТАЛ ДО ОКУПАЕМОСТИ
# ════════════════════════════════════════════════════════════════════
s = slide(); content_header(s, "Капитал", "Сколько нужно до окупаемости")
capstats=[("9,5–16,5 млн ₽","капитал до окупаемости (подушка 18–20 млн)"),("М15","операционная безубыточность"),("М23","полный возврат вложенного")]
cx=0.7
for big,small in capstats:
    box(s, cx, 1.95, 3.85, 1.5, fill=WHITE, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    text(s, cx+0.25, 2.12, 3.4, 0.75, big, size=23, color=STEEL, bold=True, font=HEAD)
    text(s, cx+0.25, 2.85, 3.4, 0.5, small, size=11.5, color=MUTED, leading=1.1)
    cx+=4.06
text(s, 0.7, 3.7, 11.9, 0.35, "ШТАТ И BURN ПО ФАЗАМ", size=12, color=STEEL, bold=True)
caprows=[("Фаза","Команда","ФОТ/мес","Полные затраты/мес"),
      ("MVP · М1–4","backend + ML/CV + 0,5 iOS","~595 тыс ₽","~626 тыс ₽"),
      ("Пилот · М5–9","+ iOS + перформанс-маркетолог","~865 тыс ₽","~1,1–1,2 млн ₽"),
      ("Рост · М10–24","+ support/sales + дизайн","~1,07 млн ₽","1,7–2,5 млн ₽")]
ry=4.1; rh=0.6; cxs=[0.7,3.2,8.3,10.2]; cws=[2.4,5.0,1.8,2.4]
for ri,row in enumerate(caprows):
    head=ri==0; fill=INK if head else (CARD if ri%2 else WHITE)
    box(s, 0.7, ry, 11.9, rh, fill=fill, line=LINE, lw=0.75)
    for ci in range(4):
        text(s, cxs[ci]+0.12, ry, cws[ci]-0.2, rh, row[ci], size=12, color=(WHITE if head else (INK if ci==0 else MUTED)), bold=(head or ci==0), anchor=MSO_ANCHOR.MIDDLE)
    ry+=rh
text(s, 0.7, ry+0.15, 11.9, 0.55, [[("Себестоимость 1 стейджинга ~10–40 ₽, маржа ~90–95%. ",{"bold":True,"color":INK}),("Рычаги ямы: цена/рамп риелторов, режим GPU (прерываемые ВМ), тайминг найма.",{"color":MUTED})]], size=12, leading=1.2)
footer(s)

# ════════════════════════════════════════════════════════════════════
# 16 — ДОРОЖНАЯ КАРТА + GATES
# ════════════════════════════════════════════════════════════════════
s = slide(); content_header(s, "План и контрольные точки", "Риски снимаются поэтапно (gates)")
stages = [
    ("Фаза 0", "Фундамент", "Авторизация + аналитика активации. До любого платного маркетинга.", STEEL),
    ("Gate 1", "Качество", "Пилот 20 риелторов: «вставил в реальный листинг» ≥ 60%.", GOLD),
    ("Gate 2", "Несущий чек", "Первый застройщик конвертируется в первом квартале пилота.", GOLD),
    ("Рост", "Масштаб", "Стейджинг → дизайнеры → фабрики/комиссия (со спросом) → Казахстан.", KEEP),
]
cw = 2.95; gap = 0.13; x0 = 0.7; y0 = 2.4
for i, (tag, h, d, c) in enumerate(stages):
    x = x0 + i*(cw+gap)
    box(s, x, y0, cw, 3.1, fill=WHITE, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    box(s, x, y0, cw, 0.7, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    box(s, x, y0+0.35, cw, 0.35, fill=c)
    text(s, x, y0, cw, 0.7, tag, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD)
    text(s, x+0.25, y0+0.9, cw-0.5, 0.5, h, size=15, color=INK, bold=True, font=HEAD)
    text(s, x+0.25, y0+1.45, cw-0.5, 1.5, d, size=12, color=MUTED, leading=1.25)
    if i < 3:
        rarrow(s, x+cw-0.03, y0+1.35, w=0.2, h=0.4, color=INK)
text(s, 0.7, 5.85, 11.9, 0.6,
     "Каждый этап — проверяемый критерий: не прошли → корректируем курс, а не жжём капитал вслепую.",
     size=12.5, color=INK, bold=True)
footer(s, 16)

# ════════════════════════════════════════════════════════════════════
# 17 — РИСКИ + СНЯТИЕ
# ════════════════════════════════════════════════════════════════════
s = slide(); content_header(s, "Риски", "Главные риски и как мы их снимаем")
risks = [
    ("Авито / Циан встроят AI-меблировку", "Институциональные договоры с застройщиками и агентствами — канал, недоступный платформе листинга; скорость; движок-ров."),
    ("Качество AI-генерации для листинга", "Движок как валидатор геометрии + сохранение реальной комнаты (inpainting); жёсткий quality-gate до маркетинга."),
    ("Узкий рынок одного сегмента", "Расширение на соседние B2B (стейджинг → застройщики → дизайнеры → Казахстан) на той же кодовой базе."),
]
yy = 2.2
for h, d in risks:
    box(s, 0.7, yy, 11.9, 1.35, fill=WHITE, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    text(s, 1.0, yy+0.2, 4.4, 1.0,
         [[("Риск", {"size": 10, "bold": True, "color": KILL})],
          [(h, {"size": 14, "bold": True, "color": INK, "font": HEAD})]], leading=1.15)
    box(s, 5.55, yy+0.2, 0.04, 0.95, fill=LINE)
    text(s, 5.85, yy+0.2, 6.5, 1.0,
         [[("Снятие", {"size": 10, "bold": True, "color": KEEP})],
          [(d, {"size": 12.5, "color": MUTED})]], leading=1.2)
    yy += 1.5
footer(s, 17)

# ════════════════════════════════════════════════════════════════════
# 18 — ИТОГ (тёмный)
# ════════════════════════════════════════════════════════════════════
s = slide(); bg(s, INK)
box(s, 0.9, 1.05, 0.9, 0.12, fill=GOLD)
kicker(s, 0.9, 1.3, "Итог")
text(s, 0.9, 1.7, 11.5, 2.2,
     [[("Будущее проекта решает не реализм AR,", {"color": WHITE})],
      [("а ", {"color": WHITE}), ("дешёвый повторяющийся спрос", {"color": GOLD}),
       (".", {"color": WHITE})]],
     size=34, bold=True, font=HEAD, leading=1.1)
text(s, 0.9, 4.0, 11.3, 1.4,
     "Мы берём деньги с профессионала (риелтор, застройщик, дизайнер), используем уже "
     "написанный код и единственный неклонируемый актив — движок расстановки. Каждый этап "
     "защищён проверяемым критерием.",
     size=15.5, color=ICE, leading=1.35)
chip(s, 0.95, 5.7, "Дальше: пилот стейджинга + identity/аналитика", STEEL, w=5.6, h=0.55)
text(s, 0.9, 6.85, 11, 0.3, "Конфиденциально · цифры — консервативные оценки, не обещания результата.",
     size=10, color=RGBColor(0x9F,0xAD,0xC2))

prs.save(OUT)
print("OK:", OUT, "·", os.path.getsize(OUT), "bytes ·", len(prs.slides._sldIdLst), "slides")
